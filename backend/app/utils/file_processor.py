"""
File upload and management utilities for document processing and vector database management
"""

import os
import logging
import hashlib
from pathlib import Path
from typing import List, Dict, Any, Optional
import asyncio
from datetime import datetime
import json
import uuid

# Document processing imports
import PyPDF2
import docx
from pptx import Presentation
import pytesseract
from pdf2image import convert_from_path
from PIL import Image
import fitz  # PyMuPDF for better PDF handling

# OpenAI for text processing and vision
import openai
from app.config import OPENAI_API_KEY

# Setup logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configure OpenAI
client = openai.OpenAI(api_key=OPENAI_API_KEY)

# File processing constants
SUPPORTED_EXTENSIONS = {
    '.pdf': 'PDF Document',
    '.docx': 'Word Document', 
    '.doc': 'Word Document (Legacy)',
    '.txt': 'Text File',
    '.md': 'Markdown File',
    '.pptx': 'PowerPoint Presentation',
    '.ppt': 'PowerPoint Presentation (Legacy)',
    '.jpg': 'JPEG Image',
    '.jpeg': 'JPEG Image',
    '.png': 'PNG Image',
    '.gif': 'GIF Image',
    '.bmp': 'Bitmap Image'
}

UPLOAD_DIR = Path("app/uploads/documents")
PROCESSED_DIR = Path("app/uploads/processed")
CACHE_DIR = Path("app/uploads/cache_documents")

# Ensure directories exist
for directory in [UPLOAD_DIR, PROCESSED_DIR, CACHE_DIR]:
    directory.mkdir(parents=True, exist_ok=True)

class FileProcessor:
    """Main class for processing uploaded files"""
    
    def __init__(self):
        self.file_metadata = {}
    
    async def process_file(self, file_path: Path, original_filename: str) -> Dict[str, Any]:
        """
        Process an uploaded file and extract text content
        
        Args:
            file_path: Path to the uploaded file
            original_filename: Original name of the file
            
        Returns:
            Dictionary containing extracted text and metadata
        """
        try:
            file_extension = file_path.suffix.lower()
            
            # Generate unique file ID
            file_id = str(uuid.uuid4())
            
            # Create metadata
            metadata = {
                "file_id": file_id,
                "original_filename": original_filename,
                "file_extension": file_extension,
                "file_size": file_path.stat().st_size,
                "upload_timestamp": datetime.now().isoformat(),
                "file_type": SUPPORTED_EXTENSIONS.get(file_extension, "Unknown"),
                "processing_status": "processing"
            }
            
            # Extract text based on file type
            extracted_text = ""
            processing_notes = []
            
            if file_extension == '.pdf':
                extracted_text, processing_notes = await self._process_pdf(file_path)
            elif file_extension in ['.docx', '.doc']:
                extracted_text, processing_notes = await self._process_word(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                extracted_text, processing_notes = await self._process_powerpoint(file_path)
            elif file_extension in ['.txt', '.md']:
                extracted_text, processing_notes = await self._process_text(file_path)
            elif file_extension in ['.jpg', '.jpeg', '.png', '.gif', '.bmp']:
                extracted_text, processing_notes = await self._process_image(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
            
            # Enhance text with LLM processing
            enhanced_text = await self._enhance_with_llm(extracted_text, metadata)
            
            # Save processed text file
            processed_file_path = PROCESSED_DIR / f"{file_id}_processed.txt"
            with open(processed_file_path, 'w', encoding='utf-8') as f:
                f.write(f"# Processed Document: {original_filename}\n")
                f.write(f"# File ID: {file_id}\n")
                f.write(f"# Processed on: {datetime.now().isoformat()}\n")
                f.write(f"# File Type: {metadata['file_type']}\n")
                f.write(f"# Original Size: {metadata['file_size']} bytes\n\n")
                
                if processing_notes:
                    f.write("## Processing Notes:\n")
                    for note in processing_notes:
                        f.write(f"- {note}\n")
                    f.write("\n")
                
                f.write("## Extracted Content:\n\n")
                f.write(enhanced_text)
            
            # Update metadata
            metadata.update({
                "processing_status": "completed",
                "processed_file_path": str(processed_file_path),
                "extracted_text_length": len(enhanced_text),
                "processing_notes": processing_notes,
                "completion_timestamp": datetime.now().isoformat()
            })
            
            # Save metadata
            metadata_path = PROCESSED_DIR / f"{file_id}_metadata.json"
            with open(metadata_path, 'w', encoding='utf-8') as f:
                json.dump(metadata, f, indent=2, ensure_ascii=False)
            
            return {
                "success": True,
                "file_id": file_id,
                "metadata": metadata,
                "extracted_text": enhanced_text,
                "processing_notes": processing_notes
            }
            
        except Exception as e:
            logger.error(f"Error processing file {original_filename}: {str(e)}")
            return {
                "success": False,
                "error": str(e),
                "file_id": None
            }
    
    async def _process_pdf(self, file_path: Path) -> tuple[str, List[str]]:
        """Process PDF files with OCR fallback"""
        text_content = []
        notes = []
        
        try:
            # Method 1: Try PyMuPDF first (better for modern PDFs)
            doc = fitz.open(file_path)
            pdf_text = ""
            
            for page_num in range(len(doc)):
                page = doc.load_page(page_num)
                page_text = page.get_text()
                
                if page_text.strip():
                    pdf_text += f"\n\n--- Page {page_num + 1} ---\n\n"
                    pdf_text += page_text
                else:
                    # If no text found, try OCR on this page
                    notes.append(f"Page {page_num + 1}: No extractable text found, attempting OCR")
                    
                    # Convert page to image for OCR
                    pix = page.get_pixmap()
                    img_data = pix.tobytes("png")
                    
                    # Use OpenAI Vision for OCR
                    ocr_text = await self._ocr_with_vision(img_data, f"PDF page {page_num + 1}")
                    if ocr_text:
                        pdf_text += f"\n\n--- Page {page_num + 1} (OCR) ---\n\n"
                        pdf_text += ocr_text
                        notes.append(f"Page {page_num + 1}: Successfully extracted text using OCR")
            
            doc.close()
            text_content.append(pdf_text)
            
        except Exception as e:
            notes.append(f"Error with PyMuPDF: {str(e)}")
            
            # Fallback: Try converting to images and OCR
            try:
                pages = convert_from_path(file_path)
                for i, page in enumerate(pages):
                    # Convert PIL image to bytes
                    import io
                    img_bytes = io.BytesIO()
                    page.save(img_bytes, format='PNG')
                    img_bytes = img_bytes.getvalue()
                    
                    ocr_text = await self._ocr_with_vision(img_bytes, f"PDF page {i + 1}")
                    if ocr_text:
                        text_content.append(f"\n\n--- Page {i + 1} (Vision OCR) ---\n\n{ocr_text}")
                        
                notes.append("Processed PDF using image conversion and Vision OCR")
                
            except Exception as e2:
                notes.append(f"Error with image conversion OCR: {str(e2)}")
                return "", notes
        
        return "\n".join(text_content), notes
    
    async def _process_word(self, file_path: Path) -> tuple[str, List[str]]:
        """Process Word documents"""
        notes = []
        
        try:
            doc = docx.Document(file_path)
            text_content = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_content.append(para.text)
            
            # Extract tables
            table_count = 0
            for table in doc.tables:
                table_count += 1
                text_content.append(f"\n\n--- Table {table_count} ---\n")
                
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        row_text.append(cell.text.strip())
                    text_content.append(" | ".join(row_text))
            
            if table_count > 0:
                notes.append(f"Extracted {table_count} tables from document")
            
            notes.append("Successfully extracted text from Word document")
            return "\n\n".join(text_content), notes
            
        except Exception as e:
            notes.append(f"Error processing Word document: {str(e)}")
            return "", notes
    
    async def _process_powerpoint(self, file_path: Path) -> tuple[str, List[str]]:
        """Process PowerPoint presentations"""
        notes = []
        
        try:
            prs = Presentation(file_path)
            text_content = []
            
            slide_count = 0
            for slide in prs.slides:
                slide_count += 1
                slide_text = f"\n\n--- Slide {slide_count} ---\n\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                text_content.append(slide_text)
            
            notes.append(f"Processed {slide_count} slides from PowerPoint presentation")
            return "\n".join(text_content), notes
            
        except Exception as e:
            notes.append(f"Error processing PowerPoint: {str(e)}")
            return "", notes
    
    async def _process_text(self, file_path: Path) -> tuple[str, List[str]]:
        """Process text and markdown files"""
        notes = []
        
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            notes.append("Successfully read text file")
            return content, notes
            
        except UnicodeDecodeError:
            # Try different encodings
            for encoding in ['latin-1', 'cp1252', 'iso-8859-1']:
                try:
                    with open(file_path, 'r', encoding=encoding) as f:
                        content = f.read()
                    notes.append(f"Successfully read text file using {encoding} encoding")
                    return content, notes
                except UnicodeDecodeError:
                    continue
            
            notes.append("Error: Could not decode text file with any encoding")
            return "", notes
        
        except Exception as e:
            notes.append(f"Error processing text file: {str(e)}")
            return "", notes
    
    async def _process_image(self, file_path: Path) -> tuple[str, List[str]]:
        """Process image files using Vision AI"""
        notes = []
        
        try:
            # Read image file
            with open(file_path, 'rb') as f:
                image_data = f.read()
            
            # Use OpenAI Vision for text extraction
            extracted_text = await self._ocr_with_vision(image_data, "image file")
            
            if extracted_text:
                notes.append("Successfully extracted text from image using OpenAI Vision")
                return extracted_text, notes
            else:
                notes.append("No text found in image")
                return "", notes
                
        except Exception as e:
            notes.append(f"Error processing image: {str(e)}")
            return "", notes
    
    async def _ocr_with_vision(self, image_data: bytes, description: str) -> str:
        """Use OpenAI Vision API for OCR"""
        try:
            import base64
            
            # Convert image to base64
            image_base64 = base64.b64encode(image_data).decode('utf-8')
            
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "text",
                                "text": f"Please extract all text content from this {description}. Include any text that appears in images, charts, diagrams, or tables. Provide a detailed description of any visual elements that contain information but no readable text. Maintain the structure and formatting as much as possible."
                            },
                            {
                                "type": "image_url",
                                "image_url": {
                                    "url": f"data:image/png;base64,{image_base64}"
                                }
                            }
                        ]
                    }
                ],
                max_tokens=2000
            )
            
            return response.choices[0].message.content
            
        except Exception as e:
            logger.error(f"Error with OpenAI Vision OCR: {str(e)}")
            return ""
    
    async def _enhance_with_llm(self, extracted_text: str, metadata: Dict[str, Any]) -> str:
        """Enhance extracted text with LLM processing"""
        try:
            if not extracted_text.strip():
                return extracted_text
            
            response = client.chat.completions.create(
                model="gpt-4o",
                messages=[
                    {
                        "role": "system",
                        "content": """You are a document processing assistant. Your task is to clean up and enhance extracted text from documents while preserving all important information. 

Please:
1. Fix any OCR errors or formatting issues
2. Improve readability while maintaining the original structure
3. Add section headers where appropriate
4. Preserve all data, numbers, and specific details
5. Add context for any unclear references
6. Maintain the professional tone of the original document"""
                    },
                    {
                        "role": "user",
                        "content": f"Please clean up and enhance this extracted text from a {metadata.get('file_type', 'document')}:\n\n{extracted_text}"
                    }
                ],
                max_tokens=4000
            )
            
            enhanced_text = response.choices[0].message.content
            return enhanced_text
            
        except Exception as e:
            logger.error(f"Error enhancing text with LLM: {str(e)}")
            return extracted_text  # Return original text if enhancement fails

# Utility functions
def get_file_hash(file_path: Path) -> str:
    """Generate SHA-256 hash of file"""
    hash_sha256 = hashlib.sha256()
    with open(file_path, "rb") as f:
        for chunk in iter(lambda: f.read(4096), b""):
            hash_sha256.update(chunk)
    return hash_sha256.hexdigest()

def is_supported_file(filename: str) -> bool:
    """Check if file extension is supported"""
    extension = Path(filename).suffix.lower()
    return extension in SUPPORTED_EXTENSIONS

def get_file_type(filename: str) -> str:
    """Get human-readable file type"""
    extension = Path(filename).suffix.lower()
    return SUPPORTED_EXTENSIONS.get(extension, "Unknown")

# Initialize processor
file_processor = FileProcessor()
